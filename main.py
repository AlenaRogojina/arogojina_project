import time
import colorsys
import tkinter as tk
import pyperclip
from tkinter import filedialog
from tkinter import colorchooser
from PIL import Image, ImageTk, ImageDraw
from PIL import Image, ImageDraw
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

class Colors_app(tk.Tk):
    def __init__(self):
        super().__init__() # Call the __init__() method of the parent class tk.Tk
        self.title("Color and Template Generator") # Set the title of the window

        self.canvas = tk.Canvas(self, width=600, height=400) #define canvas
        self.canvas.grid(row=2, column=0, columnspan=4, padx=10, pady=10)

        self.current_rgb = (None,None,None) #so each method of getting starting RGB values sets value and can be used by multiple methods. Changed to None in case color they are interested in is black (0,0,0)
        self.palette_rgb1 = (None,None,None) #sets complementary or first analogous/triadic color that can be accessed by all methods (including poster and pres)
        self.palette_rgb2 = (None,None,None) #sets second analogous/traidic color

        self.pickbtn = tk.Button(self, text="Choose a Color", command=self.colorpicker) #use colorpicker function
        self.pickbtn.grid(row=0, column=0, columnspan=2, padx=10, pady=10, sticky='we')
        
        self.imagebtn = tk.Button(self, text="Open Image", command=self.open_image) #use open image function
        self.imagebtn.grid(row=0, column=2, padx=10, pady=10, sticky='we')

        self.canvas.bind("<Motion>", self.track_cursor) #use track cursor function
        self.cursor_pos = (0, 0) #default position 0,0
        self.img_data = None #sets image to none displayed

        self.canvas.bind("<Button-1>", self.get_color) #use get color function on click

        # Create a label for the color swatch
        self.color_label = tk.Label(self.canvas, text="", bg="white", fg="black", borderwidth=1, relief="solid")
        self.color_label.place_forget()  # Hide initially

        self.randombtn = tk.Button(self, text="Random Color!", command=self.random_color) #choose a random color
        self.randombtn.grid(row=0, column=3, padx=10, pady=10, sticky='we')

        self.palettelbl = tk.Label(self, text="Choose a Palette Type:") #Label for palette type
        self.palettelbl.grid(row=1, column=0, padx=10, pady=10, sticky='we')

        self.selected = tk.StringVar(None) #creating radio buttons and tying clicking them to functions
        self.r1 = tk.Radiobutton(self, text='Complementary', value='Complementary', variable=self.selected, command=self.complementary)
        self.r1.grid(row=1, column=1, padx=10, pady=10, sticky='we')
        self.r2 = tk.Radiobutton(self, text='Analogous', value='Analogous', variable=self.selected, command=self.analogous)
        self.r2.grid(row=1, column=2, padx=10, pady=10, sticky='we')
        self.r3 = tk.Radiobutton(self, text='Triadic', value='Triadic', variable=self.selected, command=self.triadic)
        self.r3.grid(row=1, column=3, padx=10, pady=10, sticky='we')
        self.selected.set('None') #starts with none of the radio buttons selected

        self.values2 = tk.Label(self, text = "") #Label for displaying RGB values
        self.values2.grid(row=3, column=0, columnspan=4, padx=15, pady=10, sticky='w')

        self.values3 = tk.Label(self, text = "") #Label for displaying RGB values
        self.values3.place(relx = 0.5, 
                   rely = 0.89,
                   anchor = 'center')

        self.values1 = tk.Label(self, text = "") #Label for displaying RGB values
        self.values1.grid(row=3, column=2, columnspan=4, padx=20, pady=10, sticky='e')
        
        self.posterbtn = tk.Button(self, text="Download Poster Template (.pptx)", command=self.poster) #use poster function
        self.posterbtn.grid(row=4, column=0, columnspan=2, padx=10, pady=10, sticky='we')

        self.presbtn = tk.Button(self, text="Download Presentation Template (.pptx)", command=self.pres) #use presentation (pres) function
        self.presbtn.grid(row=4, column=2, columnspan=2, padx=10, pady=10, sticky='we')
    
    def rgb_to_hsv(self, r,g,b): 
        """
        Converts entered RGB values to HSV values.
        """
        return colorsys.rgb_to_hsv(r/255.0, g/255.0, b/255.0)
    
    def hsv_to_rgb(self, h, s, v):
        """
        Converts entered HSV values back to RGB values after modifying them (with analogous and triadic functions).
        """
        r, g, b = colorsys.hsv_to_rgb(h, s, v)
        return int(r*255), int(g*255), int(b*255)

    def setimage(self): 
        """
        Sets each generated image (uploaded image, color blocks) as the display image, 
        called in random color, uploaded image file, and palette functions.
        """
        self.canvas.delete("") #deletes what's diplayed on canvas currently
        self.canvas.create_image(0, 0, anchor="nw", image=self.tk_img, tags="displayed_img") #creates image based off of what is entered above inside function
        self.canvas.image = self.tk_img
        self.canvas.img_data = self.img #sets image data (in memory rather than displayed like canvas)
        self.canvas.mode = "color"

    #color picker code generated by copilot with comments by me. Also reformated to use init and self instead of root and put functions inside app
    def colorpicker(self):
        """
        Opens a color picker window, where users can select an existing color or use 
        slider and/or RGB/HSV values to create a custom one.
        """
        self.color_label.place_forget()
        self.img_data = None #so the color swatch on hover stops
        self.canvas.delete("") #deletes what's diplayed on canvas currently
        #global self.current_rgb #sets current_rgb to global so multiple functions can set/access it
        self.selected.set(None) #sets radio buttons to blank
        rgbwithhex = colorchooser.askcolor(title="Choose a color") #uses built-in color chooser, gives it a title
        self.current_rgb = rgbwithhex[0] #returns a tuple with rgb and hexcode, so only want the rgb portion
        print("Selected color:", self.current_rgb)
        self.img = Image.new('RGB', (600,600), (self.current_rgb)) #create new image with original rgb as bg
        self.tk_img = ImageTk.PhotoImage(self.img)
        self.values1.config(text = "RGB: " + str(self.current_rgb)) #Label for displaying RGB values
        self.values2.config(text = "")
        self.values3.config(text = "")

        self.setimage()
    
    def open_image(self):
        """
        Opens image file on button press.
        """
        self.color_label.place_forget()
        self.selected.set(None) #sets radio buttons to blank
        file_path = filedialog.askopenfilename() #Open file
        if file_path:
            self.img = Image.open(file_path) #The chosen file
            self.img.thumbnail((600, 400))  # Resize for display
            self.img_data = self.img.convert("RGB") #sets image data for whole class to currently display opened file
            self.tk_img = ImageTk.PhotoImage(self.img)

            self.setimage()

    def track_cursor(self, event):
        """
        Tracks mouse position to use for label that shows RGB values when hovering over uploaded image.
        """
        self.color_label.place_forget()
        self.cursor_pos = (event.x, event.y) #gets up to date x and y coordinates of cursor as it moves
        if self.img_data:
            try:
                rgb = self.img_data.getpixel((self.cursor_pos)) #gets RGB of where cursor is on the image
                hex_color = "#%02x%02x%02x" % rgb #convert to hex because next line needs hex not RGB
                self.color_label.config(bg=hex_color, text=rgb) #shows color and its RGB value
                self.color_label.place(x=event.x + 10, y=event.y + 10) #places a bit to the side of the cursor
            except IndexError:
                self.color_label.place_forget() #does not display outside of image

    def get_color(self, event):
        """
        Gets pixel color that mouse is hovering over.
        """
        #global current_rgb
        self.color_label.place_forget()
        x, y = event.x, event.y #same as above
        if self.img_data:
            try:
                rgb = self.img_data.getpixel((x, y)) #
                self.current_rgb = rgb
                print(f"Cursor at ({x}, {y}) - Color: {self.current_rgb}")
                self.values1.config(text = "RGB: " + str(self.current_rgb)) #Label for displaying RGB values
                self.values2.config(text = "")
                self.values3.config(text = "")
            except IndexError:
                print("Cursor is outside the image bounds.")

    def random_color(self):
        """
        Generates a random starting color.
        """
        self.canvas.delete("") #deletes what's diplayed on canvas currently
        self.img_data = None #so the color swatch on hover stops
        self.color_label.place_forget()
        #global current_rgb
        self.selected.set(None) #set radio buttons to blank
        self.current_rgb = (random.randint(0,255), random.randint(0,255), random.randint(0,255))
        print("Generated color: ", self.current_rgb)

        self.img = Image.new('RGB', (600,600), (self.current_rgb)) #create new image with original rgb as bg
        self.tk_img = ImageTk.PhotoImage(self.img)
        self.values1.config(text = "RGB: " + str(self.current_rgb)) #Label for displaying RGB values
        self.values2.config(text = "")
        self.values3.config(text = "")

        self.setimage()
    
    def complementary(self):
        """
        Generates a complementary color for the starting color and displays two side by side.
        """
        self.palette_rgb2 = (None,None,None) #was running into issue where if analogous was chosen first, template created, then complementary chosen and template created, kept third color in the powerpoint
        self.img_data = None #so the color swatch on hover stops
        self.color_label.place_forget()
        #global current_rgb
        r, g, b = self.current_rgb
        h, s, v = self.rgb_to_hsv(r, g, b)
        h = (h + 0.5) % 1.0  # 180 degrees around the hue circle
        comp_r, comp_g, comp_b = self.hsv_to_rgb(h, s, v)

        self.palette_rgb1 = (comp_r, comp_g, comp_b) #get new values and set as palette color 1

        print(self.palette_rgb1)
        self.values2.config(text = "RGB: " + str(self.palette_rgb1)) #Label for displaying RGB values
        self.values3.config(text = "")

        self.img = Image.new('RGB', (600,600), (self.current_rgb)) #create new image with original rgb as background
        self.draw = ImageDraw.Draw(self.img) #draw image
        self.draw.rectangle((0, 0, 300, 402), fill=(self.palette_rgb1), outline=(self.palette_rgb1))

        self.tk_img = ImageTk.PhotoImage(self.img)
        self.setimage()

    def analogous(self):
        """
        Generates two analogous colors for the starting color and displays three side by side.
        """
        self.img_data = None #so the color swatch on hover stops
        self.color_label.place_forget()
        #global current_rgb
        r, g, b = self.current_rgb
        h, s, v = self.rgb_to_hsv(r, g, b)
        h1 = (h + 0.083) % 1.0  # 30 degrees around the hue circle
        h2 = (h - 0.083) % 1.0  # 30 degrees around the hue circle in opposite direction
        self.palette_rgb1 = self.hsv_to_rgb(h1, s, v)
        self.palette_rgb2 = self.hsv_to_rgb(h2, s, v)

        print([self.palette_rgb1, self.palette_rgb2])
        self.values2.config(text = "RGB: " + str(self.palette_rgb1)) #Label for displaying RGB values
        self.values3.config(text = "RGB: " + str(self.palette_rgb2))

        self.img = Image.new('RGB', (600, 600), (r, g, b)) #create new image
        self.draw = ImageDraw.Draw(self.img) #draw
        self.draw.rectangle((0, 0, 600//3, 402), fill=(self.palette_rgb1)) #draw first analogous color rectangle
        self.draw.rectangle((600//3, 0, (600//3 + 600//3), 402), fill=(self.palette_rgb2)) #second

        self.tk_img = ImageTk.PhotoImage(self.img)
        self.setimage()
    
    def triadic(self):
        """
        Generates two triadic colors for the starting color and displays three side by side.
        """
        self.img_data = None #so the color swatch on hover stops
        self.color_label.place_forget()
        #global current_rgb
        r, g, b = self.current_rgb
        h, s, v = self.rgb_to_hsv(r, g, b)
        h1 = (h + 0.333) % 1.0  # 30 degrees around the hue circle
        h2 = (h - 0.333) % 1.0  # 30 degrees around the hue circle in opposite direction
        self.palette_rgb1 = self.hsv_to_rgb(h1, s, v)
        self.palette_rgb2 = self.hsv_to_rgb(h2, s, v)

        print([self.palette_rgb1, self.palette_rgb2])
        self.values2.config(text = "RGB: " + str(self.palette_rgb1)) #Label for displaying RGB values
        self.values3.config(text = "RGB: " + str(self.palette_rgb2))

        self.img = Image.new('RGB', (600, 600), (r, g, b)) #create new image
        self.draw = ImageDraw.Draw(self.img) #draw
        self.draw.rectangle((0, 0, 600//3, 402), fill=(self.palette_rgb1)) #draw first triadic color rectangle
        self.draw.rectangle((600//3, 0, (600//3 + 600//3), 402), fill=(self.palette_rgb2)) #second

        self.tk_img = ImageTk.PhotoImage(self.img)
        self.setimage()

    def force_bullet(self,paragraph, char="•"):
        """
        Ensure the paragraph in a textbox shows a bullet.
        """
        pPr = paragraph._p.get_or_add_pPr()
        # Remove explicit 'no bullet'
        for el in pPr.findall(qn('a:buNone')):
            pPr.remove(el)
        # Add bullet char if not present
        if not pPr.findall(qn('a:buChar')):
            buChar = OxmlElement('a:buChar')
            buChar.set('char', char)
            pPr.append(buChar)

    def poster(self):
        """
        Creates a poster pptx template, with intro, methods, results, discussion, take home point, and a contact/references section.
        """
        prs = Presentation()
        prs.slide_width = Inches(36)
        prs.slide_height = Inches(24)
        # Use the blank slide layout (index 6) for a clean slide
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        # Background color (original color--pale blue)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(self.current_rgb[0], self.current_rgb[1], self.current_rgb[2])   #just using self.current_rgb doesn't work, which is stupid. I guess I could revise to make it a list

        # Title
        title_box = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(36), Inches(3.3)) #create background box for title
        title_box.fill.solid() #set box color
        title_box.fill.fore_color.rgb = RGBColor(255, 255, 255) #set text color for title
        p = title_box.text_frame.paragraphs[0] #create text frame for title
        run = p.add_run()
        run.text = "Title and Names Here"
        run.font.size = Pt(72)
        run.font.bold = True
        if self.palette_rgb2 == (None,None,None): #if complementary (no third color)
            run.font.color.rgb = RGBColor(self.palette_rgb1[0], self.palette_rgb1[1], self.palette_rgb1[2])
        else:
            run.font.color.rgb = RGBColor(self.palette_rgb2[0], self.palette_rgb2[1], self.palette_rgb2[2])
        p.alignment = PP_ALIGN.CENTER

        # Introduction
        intro = slide.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(11), Inches(0.7)) #create box for "Introduction"
        intro.fill.solid()
        intro.fill.fore_color.rgb = RGBColor(self.palette_rgb1[0], self.palette_rgb1[1], self.palette_rgb1[2]) #use first palette color for all palette types
        tf = intro.text_frame.paragraphs[0]
        runintro = tf.add_run()
        runintro.text = "Introduction"
        runintro.font.size = Pt(40)
        runintro.font.bold = True
        runintro.font.color.rgb = RGBColor(255, 255, 255)
        tf.alignment = PP_ALIGN.CENTER

        tintro = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(11), Inches(8.1)) #create box for text of intro
        tintro.fill.solid()
        tintro.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tfintro = tintro.text_frame
        items = ["background", "research question", "hypotheses"]
        for i, txt in enumerate(items):
            para = tfintro.paragraphs[0] if i == 0 else tfintro.add_paragraph()
            para.text = txt
            para.level = 0
            self.force_bullet(para)
            para.font.size = Pt(25)
            para.font.color.rgb = RGBColor(0, 0, 0)

        # Methods
        methods = slide.shapes.add_textbox(Inches(0.7), Inches(13.1), Inches(11), Inches(0.7))
        methods.fill.solid()
        methods.fill.fore_color.rgb = RGBColor(self.palette_rgb1[0], self.palette_rgb1[1], self.palette_rgb1[2])
        tfm = methods.text_frame.paragraphs[0]
        runmethods = tfm.add_run()
        runmethods.text = "Methods"
        runmethods.font.size = Pt(40)
        runmethods.font.bold = True
        runmethods.font.color.rgb = RGBColor(255, 255, 255)
        tfm.alignment = PP_ALIGN.CENTER

        tmethods = slide.shapes.add_textbox(Inches(0.7), Inches(13.8), Inches(11), Inches(9.4))
        tmethods.fill.solid()
        tmethods.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tfmethods = tmethods.text_frame
        items = ["procedure", "participants", "materials"]
        for i, txt in enumerate(items):
            para = tfmethods.paragraphs[0] if i == 0 else tfmethods.add_paragraph()
            para.text = txt
            para.level = 0
            self.force_bullet(para)
            para.font.size = Pt(25)
            para.font.color.rgb = RGBColor(0, 0, 0)

        # Results
        results = slide.shapes.add_textbox(Inches(12.4), Inches(4.1), Inches(11.1), Inches(0.7))
        results.fill.solid()
        results.fill.fore_color.rgb = RGBColor(self.palette_rgb1[0], self.palette_rgb1[1], self.palette_rgb1[2])
        tfr = results.text_frame.paragraphs[0]
        runresults = tfr.add_run()
        runresults.text = "Results"
        runresults.font.size = Pt(40)
        runresults.font.bold = True
        runresults.font.color.rgb = RGBColor(255, 255, 255)
        tfr.alignment = PP_ALIGN.CENTER

        tresults = slide.shapes.add_textbox(Inches(12.4), Inches(4.8), Inches(11.1), Inches(18.4))
        tresults.fill.solid()
        tresults.fill.fore_color.rgb = RGBColor(255, 255, 255)
        #tfresults = tresults.text_frame

        if self.palette_rgb2 != (None,None,None):
            #Explaining third color
            thirdcol = slide.shapes.add_textbox(Inches(13), Inches(5.5), Inches(9), Inches(1))
            tftc = thirdcol.text_frame.paragraphs[0]
            runtc = tftc.add_run()
            runtc.text = "Use the third color for accents and highlights"
            runtc.font.size = Pt(38)
            runtc.font.bold = True
            runtc.font.color.rgb = RGBColor(self.palette_rgb2[0], self.palette_rgb2[1], self.palette_rgb2[2])
            tfr.alignment = PP_ALIGN.CENTER

        # Discussion
        disc = slide.shapes.add_textbox(Inches(24.3), Inches(4.1), Inches(11), Inches(0.7))
        disc.fill.solid()
        disc.fill.fore_color.rgb = RGBColor(self.palette_rgb1[0], self.palette_rgb1[1], self.palette_rgb1[2])
        tfd = disc.text_frame.paragraphs[0]
        rundisc = tfd.add_run()
        rundisc.text = "Discussion"
        rundisc.font.size = Pt(40)
        rundisc.font.bold = True
        rundisc.font.color.rgb = RGBColor(255, 255, 255)
        tfd.alignment = PP_ALIGN.CENTER

        tdisc = slide.shapes.add_textbox(Inches(24.3), Inches(4.8), Inches(11), Inches(12.4))
        tdisc.fill.solid()
        tdisc.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tfdisc = tdisc.text_frame
        items = ["hypotheses support", "limitations", "conclusion"]
        for i, txt in enumerate(items):
            para = tfdisc.paragraphs[0] if i == 0 else tfdisc.add_paragraph()
            para.text = txt
            para.level = 0
            self.force_bullet(para)
            para.font.size = Pt(25)
            para.font.color.rgb = RGBColor(0, 0, 0)

        # Take Home Point
        thp = slide.shapes.add_textbox(Inches(24.3), Inches(17.4), Inches(11), Inches(0.7))
        thp.fill.solid()
        if self.palette_rgb2 == (None,None,None): #if complementary (no third color)
            thp.fill.fore_color.rgb = RGBColor(self.palette_rgb1[0], self.palette_rgb1[1], self.palette_rgb1[2])
        else:
            thp.fill.fore_color.rgb = RGBColor(self.palette_rgb2[0], self.palette_rgb2[1], self.palette_rgb2[2])
        tft = thp.text_frame.paragraphs[0]
        runt = tft.add_run()
        runt.text = "Take Home Point/Main Finding"
        runt.font.size = Pt(40)
        runt.font.bold = True
        runt.font.color.rgb = RGBColor(255, 255, 255)
        tft.alignment = PP_ALIGN.CENTER

        tthp = slide.shapes.add_textbox(Inches(24.3), Inches(18.1), Inches(11), Inches(3))
        tthp.fill.solid()
        tthp.fill.fore_color.rgb = RGBColor(255, 255, 255)
        #tfdisc = tdisc.text_frame

        # Contact and References
        cr = slide.shapes.add_textbox(Inches(24.3), Inches(21.4), Inches(11), Inches(1.72))
        cr.fill.solid()
        cr.fill.fore_color.rgb = RGBColor(self.palette_rgb1[0], self.palette_rgb1[1], self.palette_rgb1[2])
        tfcr = cr.text_frame.paragraphs[0]
        runcr = tfcr.add_run()
        runcr.text = "Contact and References (QR code)"
        runcr.font.size = Pt(40)
        runcr.font.bold = True
        runcr.font.color.rgb = RGBColor(255, 255, 255)
        tfcr.alignment = PP_ALIGN.LEFT

        timestr = time.strftime("%Y%m%d-%H%M%S")
        prs.save("poster" + timestr + ".pptx")

    def pres(self):
        """
        Creates a presentation pptx template with title slide and one information slide.
        """
        prs = Presentation()
        prs.slide_width = Inches(12)
        prs.slide_height = Inches(6)
        # Use the blank slide layout (index 6) for a clean slide
        # other options: 0: Title Slide 1: Title and Content 2: Section Header 3: Two Content 4: Comparison 5: Title Only 6: Blank 7: Content with Caption 8: Picture with Caption
        blank = prs.slide_layouts[6]
        titleslide = prs.slides.add_slide(blank)

        # Background color (original color)
        bg = titleslide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(self.current_rgb[0], self.current_rgb[1], self.current_rgb[2])

        # Title
        title_box = titleslide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(10), Inches(2.5))
        p = title_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Title Here"
        run.font.size = Pt(70)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        #title line
        tline = titleslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4), Inches(12), Inches(0.3))
        tline.fill.solid()
        tline.fill.fore_color.rgb = RGBColor(self.palette_rgb1[0], self.palette_rgb1[1], self.palette_rgb1[2])
        tline.shadow.inherit = False
        tline.line.fill.background()
        
        if self.palette_rgb2 != (None,None,None): #if complementary palette chosen
            #Title Box
            tbox = titleslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.3), Inches(12), Inches(1.7))
            tbox.fill.solid()
            tbox.fill.fore_color.rgb = RGBColor(self.palette_rgb2[0], self.palette_rgb2[1], self.palette_rgb2[2])
            tbox.shadow.inherit = False
            tbox.line.fill.background()

        #Names
        name_box = titleslide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(10), Inches(2))
        name = name_box.text_frame.paragraphs[0]
        runn = name.add_run()
        runn.text = "Additional Info Here"
        runn.font.size = Pt(30)
        runn.font.bold = True
        runn.font.color.rgb = RGBColor(255, 255, 255)
        name.alignment = PP_ALIGN.CENTER

        #Second Slide--info slide
        secondslide = prs.slides.add_slide(blank)
        # Background color for first info slide (original color)
        bg2 = secondslide.background
        bg2.fill.solid()
        bg2.fill.fore_color.rgb = RGBColor(self.current_rgb[0], self.current_rgb[1], self.current_rgb[2])

        sbox = secondslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(12), Inches(1.5))
        sbox.fill.solid()
        sbox.fill.fore_color.rgb = RGBColor(self.palette_rgb1[0], self.palette_rgb1[1], self.palette_rgb1[2])
        sbox.shadow.inherit = False
        sbox.line.fill.background()

        if self.palette_rgb2 != (None,None,None):
            #Second Slide Second Line
            sbox = secondslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5), Inches(12), Inches(0.3))
            sbox.fill.solid()
            sbox.fill.fore_color.rgb = RGBColor(self.palette_rgb2[0], self.palette_rgb2[1], self.palette_rgb2[2])
            sbox.shadow.inherit = False
            sbox.line.fill.background()

        # Second Slide Title
        two_title_box = secondslide.shapes.add_textbox(Inches(0.75), Inches(0.75), Inches(10), Inches(2))
        twotitle = two_title_box.text_frame.paragraphs[0]
        runtwo = twotitle.add_run()
        runtwo.text = "Title Here"
        runtwo.font.size = Pt(40)
        runtwo.font.bold = True
        runtwo.font.color.rgb = RGBColor(255, 255, 255)
        twotitle.alignment = PP_ALIGN.LEFT

        #Second Slide Text
        ttwo = secondslide.shapes.add_textbox(Inches(0.75), Inches(1.75), Inches(10), Inches(4))
        tftwo = ttwo.text_frame
        items = ["point 1", "point 2", "point 3"]
        for i, txt in enumerate(items):
            para = tftwo.paragraphs[0] if i == 0 else tftwo.add_paragraph()
            para.text = txt
            para.level = 0
            self.force_bullet(para)
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(255,255,255)

        timestr = time.strftime("%Y%m%d-%H%M%S")
        prs.save("presentation" + timestr + ".pptx")

app = Colors_app()
app.mainloop()
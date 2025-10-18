#pip install python-pptx
from PIL import Image, ImageDraw
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Grab input color (for 1)
#sr, sg, sb = int(input("Enter the three rgb values:"))
sr = int(input("Enter the r value:"))
sg = int(input("Enter the g value:"))
sb = int(input("Enter the b value:"))


# Generate random color (for 2)
def randomcolor():
    r = random.randint(0,255)
    g = random.randint(0,255)
    b = random.randint(0,255)

    return r, g, b

randr, randg, randb = randomcolor()

# Generate new colors based on input color
def complementary(r, g, b):
    newr = 255 - r
    newg = 255 - g #get new values
    newb = 255 - b

    return newr, newg, newb

compr, compg, compb = complementary(sr, sg, sb) #get values
print(compr, compg, compb)

def analogous(r,g,b):
    newr1 = r
    newg1 = g + 31 #get new values
    if newg1 > 255:
        newg1 = g - 31 #make sure new values aren't higher than 255
    newb1 = b

    newr2 = r
    newg2 = g + 16
    if newg2 > 255:
        newg2 = g - 16
    newb2 = b

    return [newr1, newg1, newb1, newr2, newg2, newb2]

analog = analogous(sr,sg,sb) #get analogous values
print(analog)

def tetrad(r,g,b):
    newr1 = r + 71
    if newr1 > 255:
        newr1 = g - 71 #make sure new values aren't higher than 255
    newb1 = b
    newg1 = g - 46
    if newg1 < 0:
        newg1 = g + 46 #make sure new values aren't lower than 0
    newb1 = b
    newb1 = b

    newr2 = r + 78
    if newr2 > 255:
        newr2 = g - 78
    newg2 = g - 14
    if newg2 < 0:
        newg2 = g + 14
    newb2 = b - 78
    if newb2 < 0:
        newb2 = g + 78

    return [newr1, newg1, newb1, newr2, newg2, newb2]


tet = tetrad(sr, sg, sb)
print(tet)

def force_bullet(paragraph, char="•"):
    """
    Ensure the paragraph in a textbox shows a bullet.
    """
    pPr = paragraph._p.get_or_add_pPr()
    # Remove explicit 'no bullet'
    for el in pPr.findall(qn('a:buNone')):
        pPr.remove(el)
    # Add bullet char if not present
    if not pPr.findall(qn('a:buChar')):
        buChar = OxmlElement('a:buChar')
        buChar.set('char', char)
        pPr.append(buChar)

# Radio buttons for complementary, analogous, triadic
paltype = input("What kind of palette would you like? Type in 'complementary', 'analogous', or 'tetrad'.")
def palettedisplay(r,g, b, paltype):
    if paltype == "complementary":
        nr, ng, nb = complementary(r, g, b)
        im = Image.new('RGB', (500, 300), (r, g, b)) #create new image with original rgb as bg
        draw = ImageDraw.Draw(im) #draw image

        draw.rectangle((0, 0, 250, 300), fill=(nr, ng, nb), outline=(nr, ng, nb)) #add complementary color rectangle
        im.save('pillow_imagedraw_complementary2.jpg', quality=95) #save file
    elif paltype == "analogous":
        nr, ng, nb = analogous(r, g, b)
        imanalogous = Image.new('RGB', (500, 300), (r, g, b)) #create new image
        draw = ImageDraw.Draw(imanalogous) #draw

        draw.rectangle((0, 0, 500//3, 300), fill=(analog[0], analog[1], analog[2])) #draw first analogous color rectangle
        draw.rectangle((500//3, 0, (500//3 + 500//3), 300), fill=(analog[3], analog[4], analog[5])) #second
        imanalogous.save('pillow_imagedraw_analogous2.jpg', quality=95) #save file
    elif paltype == "tetrad":
        nr, ng, nb = tetrad(r, g, b)
        imtetrad = Image.new('RGB', (500, 300), (r, g, b)) #create new image
        draw = ImageDraw.Draw(imtetrad) #draw

        draw.rectangle((0, 0, 500//3, 300), fill=(tet[0], tet[1], tet[2])) #draw first analogous color rectangle
        draw.rectangle((500//3, 0, (500//3 + 500//3), 300), fill=(tet[3], tet[4], tet[5])) #second
        imtetrad.save('pillow_imagedraw_tetrad2.jpg', quality=95) #save file
    else:
        print("Please type in 'complementary', 'analogous', or 'tetrad'.")

# Output color blocks
palettedisplay(sr, sg, sb, paltype)

# Output template
temptype = input("What kind of template would you like? Type in 'presentation' or 'poster'.")
def template(r,g,b,paltype,temptype):
    if temptype == "presentation":
        compr, compg, compb = complementary(r, g, b)
        analog = analogous(r,g,b)
        tet = tetrad(r,g,b)
        prs = Presentation()
        prs.slide_width = Inches(12)
        prs.slide_height = Inches(6)
        # Use the blank slide layout (index 6) for a clean slide
        # other options: 0: Title Slide 1: Title and Content 2: Section Header 3: Two Content 4: Comparison 5: Title Only 6: Blank 7: Content with Caption 8: Picture with Caption
        blank = prs.slide_layouts[6]
        titleslide = prs.slides.add_slide(blank)

        # Background color (original color)
        bg = titleslide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

        # Title
        title_box = titleslide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(10), Inches(2.5))
        p = title_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Title Here"
        run.font.size = Pt(70)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        if paltype == "complementary":
            #Title line
            tline = titleslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4), Inches(12), Inches(0.3))
            tline.fill.solid()
            tline.fill.fore_color.rgb = RGBColor(compr, compg, compb)
            tline.shadow.inherit = False
            tline.line.fill.background()
        
        elif paltype == "analogous":
            #Title line
            tline = titleslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4), Inches(12), Inches(0.3))
            tline.fill.solid()
            tline.fill.fore_color.rgb = RGBColor(analog[3], analog[4], analog[5])
            tline.shadow.inherit = False
            tline.line.fill.background()

            #Title Box
            tbox = titleslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.3), Inches(12), Inches(1.7))
            tbox.fill.solid()
            tbox.fill.fore_color.rgb = RGBColor(analog[0], analog[1], analog[2])
            tbox.shadow.inherit = False
            tbox.line.fill.background()

        elif paltype == "tetrad":
            #Title line
            #title_line = titleslide.shapes.rectangle(Inches(0), Inches(4), Inches(12), Inches(0.5))
            tline = titleslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4), Inches(12), Inches(0.3))
            tline.fill.solid()
            tline.fill.fore_color.rgb = RGBColor(tet[3], tet[4], tet[5])
            tline.shadow.inherit = False
            tline.line.fill.background()

            #Title Box
            tbox = titleslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.3), Inches(12), Inches(1.7))
            tbox.fill.solid()
            tbox.fill.fore_color.rgb = RGBColor(tet[0], tet[1], tet[2])
            tbox.shadow.inherit = False
            tbox.line.fill.background()

        #Names
        name_box = titleslide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(10), Inches(2))
        name = name_box.text_frame.paragraphs[0]
        runn = name.add_run()
        runn.text = "Additional Info Here"
        runn.font.size = Pt(30)
        runn.font.bold = True
        runn.font.color.rgb = RGBColor(255, 255, 255)
        name.alignment = PP_ALIGN.CENTER

        #Second Slide--info slide
        secondslide = prs.slides.add_slide(blank)
        # Background color for first info slide (original color)
        bg2 = secondslide.background
        bg2.fill.solid()
        bg2.fill.fore_color.rgb = RGBColor(r, g, b)

        if paltype == "analogous":
            #Second Slide Box
            sbox = secondslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(12), Inches(1.5))
            sbox.fill.solid()
            sbox.fill.fore_color.rgb = RGBColor(analog[0], analog[1], analog[2])
            sbox.shadow.inherit = False
            sbox.line.fill.background()

            #Second Slide Second Line
            sbox = secondslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5), Inches(12), Inches(0.3))
            sbox.fill.solid()
            sbox.fill.fore_color.rgb = RGBColor(analog[3], analog[4], analog[5])
            sbox.shadow.inherit = False
            sbox.line.fill.background()

        if paltype == "tetrad":
            #Second Slide Box
            sbox = secondslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(12), Inches(1.5))
            sbox.fill.solid()
            sbox.fill.fore_color.rgb = RGBColor(tet[0], tet[1], tet[2])
            sbox.shadow.inherit = False
            sbox.line.fill.background()

            #Second Slide Second Line
            sbox = secondslide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5), Inches(12), Inches(0.3))
            sbox.fill.solid()
            sbox.fill.fore_color.rgb = RGBColor(tet[3], tet[4], tet[5])
            sbox.shadow.inherit = False
            sbox.line.fill.background()

        # Second Slide Title
        two_title_box = secondslide.shapes.add_textbox(Inches(0.75), Inches(0.75), Inches(10), Inches(2))
        twotitle = two_title_box.text_frame.paragraphs[0]
        runtwo = twotitle.add_run()
        runtwo.text = "Title Here"
        runtwo.font.size = Pt(40)
        runtwo.font.bold = True
        runtwo.font.color.rgb = RGBColor(255, 255, 255)
        twotitle.alignment = PP_ALIGN.LEFT

        #Second Slide Text
        ttwo = secondslide.shapes.add_textbox(Inches(0.75), Inches(1.75), Inches(10), Inches(4))
        tftwo = ttwo.text_frame
        items = ["point 1", "point 2", "point 3"]
        for i, txt in enumerate(items):
            para = tftwo.paragraphs[0] if i == 0 else tftwo.add_paragraph()
            para.text = txt
            para.level = 0
            force_bullet(para)
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(255,255,255)

        prs.save("prestemplate-function.pptx")

    elif temptype == "poster":
        prs = Presentation()
        prs.slide_width = Inches(36)
        prs.slide_height = Inches(24)
        # Use the blank slide layout (index 6) for a clean slide
        # other options: 0: Title Slide 1: Title and Content 2: Section Header 3: Two Content 4: Comparison 5: Title Only 6: Blank 7: Content with Caption 8: Picture with Caption
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        # Background color (original color--pale blue)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(36), Inches(3.3))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        p = title_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Title and Names Here"
        run.font.size = Pt(72)
        run.font.bold = True
        if paltype == "complementary":
            run.font.color.rgb = RGBColor(compr, compg, compb)
        elif paltype == "analogous":
            run.font.color.rgb = RGBColor(analog[3], analog[4], analog[5])
        elif paltype == "tetrad":
            run.font.color.rgb = RGBColor(tet[3], tet[4], tet[5])
        p.alignment = PP_ALIGN.CENTER

        # Introduction
        intro = slide.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(11), Inches(0.7))
        intro.fill.solid()
        if paltype == "complementary":
            intro.fill.fore_color.rgb = RGBColor(compr, compg, compb)
        elif paltype == "analogous":
            intro.fill.fore_color.rgb = RGBColor(analog[0], analog[1], analog[2])
        elif paltype == "tetrad":
            intro.fill.fore_color.rgb = RGBColor(tet[0], tet[1], tet[2])
        tf = intro.text_frame.paragraphs[0]
        runintro = tf.add_run()
        runintro.text = "Introduction"
        runintro.font.size = Pt(40)
        runintro.font.bold = True
        runintro.font.color.rgb = RGBColor(255, 255, 255)
        tf.alignment = PP_ALIGN.CENTER

        tintro = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(11), Inches(8.1))
        tintro.fill.solid()
        tintro.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tfintro = tintro.text_frame
        items = ["background", "research question", "hypotheses"]
        for i, txt in enumerate(items):
            para = tfintro.paragraphs[0] if i == 0 else tfintro.add_paragraph()
            para.text = txt
            para.level = 0
            force_bullet(para)
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(0, 0, 0)

        # Methods
        methods = slide.shapes.add_textbox(Inches(0.7), Inches(13.1), Inches(11), Inches(0.7))
        methods.fill.solid()
        if paltype == "complementary":
            methods.fill.fore_color.rgb = RGBColor(compr, compg, compb)
        elif paltype == "analogous":
            methods.fill.fore_color.rgb = RGBColor(analog[0], analog[1], analog[2])
        elif paltype == "tetrad":
            methods.fill.fore_color.rgb = RGBColor(tet[0], tet[1], tet[2])
        tfm = methods.text_frame.paragraphs[0]
        runmethods = tfm.add_run()
        runmethods.text = "Methods"
        runmethods.font.size = Pt(40)
        runmethods.font.bold = True
        runmethods.font.color.rgb = RGBColor(255, 255, 255)
        tfm.alignment = PP_ALIGN.CENTER

        tmethods = slide.shapes.add_textbox(Inches(0.7), Inches(13.8), Inches(11), Inches(9.4))
        tmethods.fill.solid()
        tmethods.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tfmethods = tmethods.text_frame
        items = ["procedure", "participants", "materials"]
        for i, txt in enumerate(items):
            para = tfmethods.paragraphs[0] if i == 0 else tfmethods.add_paragraph()
            para.text = txt
            para.level = 0
            force_bullet(para)
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(0, 0, 0)

        # Results
        results = slide.shapes.add_textbox(Inches(12.4), Inches(4.1), Inches(11.1), Inches(0.7))
        results.fill.solid()
        if paltype == "complementary":
            results.fill.fore_color.rgb = RGBColor(compr, compg, compb)
        elif paltype == "analogous":
            results.fill.fore_color.rgb = RGBColor(analog[0], analog[1], analog[2])
        elif paltype == "tetrad":
            results.fill.fore_color.rgb = RGBColor(tet[0], tet[1], tet[2])
        tfr = results.text_frame.paragraphs[0]
        runresults = tfr.add_run()
        runresults.text = "Results"
        runresults.font.size = Pt(40)
        runresults.font.bold = True
        runresults.font.color.rgb = RGBColor(255, 255, 255)
        tfr.alignment = PP_ALIGN.CENTER

        tresults = slide.shapes.add_textbox(Inches(12.4), Inches(4.8), Inches(11.1), Inches(18.4))
        tresults.fill.solid()
        tresults.fill.fore_color.rgb = RGBColor(255, 255, 255)
        #tfresults = tresults.text_frame

        if paltype == "analogous" or "complementary":
            #Explaining third color
            thirdcol = slide.shapes.add_textbox(Inches(13), Inches(5.5), Inches(9), Inches(1))
            tftc = thirdcol.text_frame.paragraphs[0]
            runtc = tftc.add_run()
            runtc.text = "Use the third color for accents and highlights"
            runtc.font.size = Pt(38)
            runtc.font.bold = True
            if paltype == "analogous":
                runtc.font.color.rgb = RGBColor(analog[3], analog[4], analog[5])
            elif paltype == "tetrad":
                runtc.font.color.rgb = RGBColor(tet[3], tet[4], tet[5])
            tfr.alignment = PP_ALIGN.CENTER

        # Discussion
        disc = slide.shapes.add_textbox(Inches(24.3), Inches(4.1), Inches(11), Inches(0.7))
        disc.fill.solid()
        if paltype == "complementary":
            disc.fill.fore_color.rgb = RGBColor(compr, compg, compb)
        elif paltype == "analogous":
            disc.fill.fore_color.rgb = RGBColor(analog[0], analog[1], analog[2])
        elif paltype == "tetrad":
            disc.fill.fore_color.rgb = RGBColor(tet[0], tet[1], tet[2])
        tfd = disc.text_frame.paragraphs[0]
        rundisc = tfd.add_run()
        rundisc.text = "Discussion"
        rundisc.font.size = Pt(40)
        rundisc.font.bold = True
        rundisc.font.color.rgb = RGBColor(255, 255, 255)
        tfd.alignment = PP_ALIGN.CENTER

        tdisc = slide.shapes.add_textbox(Inches(24.3), Inches(4.8), Inches(11), Inches(12.4))
        tdisc.fill.solid()
        tdisc.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tfdisc = tdisc.text_frame
        items = ["hypotheses support", "limitations", "conclusion"]
        for i, txt in enumerate(items):
            para = tfdisc.paragraphs[0] if i == 0 else tfdisc.add_paragraph()
            para.text = txt
            para.level = 0
            force_bullet(para)
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(0, 0, 0)

        # Take Home Point
        thp = slide.shapes.add_textbox(Inches(24.3), Inches(17.4), Inches(11), Inches(0.7))
        thp.fill.solid()
        if paltype == "complementary":
            thp.fill.fore_color.rgb = RGBColor(compr, compg, compb)
        elif paltype == "analogous":
            thp.fill.fore_color.rgb = RGBColor(analog[3], analog[4], analog[5])
        elif paltype == "tetrad":
            thp.fill.fore_color.rgb = RGBColor(tet[3], tet[4], tet[5])
        tft = thp.text_frame.paragraphs[0]
        runt = tft.add_run()
        runt.text = "Take Home Point/Main Finding"
        runt.font.size = Pt(40)
        runt.font.bold = True
        runt.font.color.rgb = RGBColor(255, 255, 255)
        tft.alignment = PP_ALIGN.CENTER

        tthp = slide.shapes.add_textbox(Inches(24.3), Inches(18.1), Inches(11), Inches(3))
        tthp.fill.solid()
        tthp.fill.fore_color.rgb = RGBColor(255, 255, 255)
        #tfdisc = tdisc.text_frame

        # Contact and References
        cr = slide.shapes.add_textbox(Inches(24.3), Inches(21.4), Inches(11), Inches(1.72))
        cr.fill.solid()
        if paltype == "complementary":
            cr.fill.fore_color.rgb = RGBColor(compr, compg, compb)
        elif paltype == "analogous":
            cr.fill.fore_color.rgb = RGBColor(analog[0], analog[1], analog[2])
        elif paltype == "tetrad":
            cr.fill.fore_color.rgb = RGBColor(tet[0], tet[1], tet[2])
        tfcr = cr.text_frame.paragraphs[0]
        runcr = tfcr.add_run()
        runcr.text = "Contact and References (QR code)"
        runcr.font.size = Pt(40)
        runcr.font.bold = True
        runcr.font.color.rgb = RGBColor(255, 255, 255)
        tfcr.alignment = PP_ALIGN.LEFT

        prs.save("postertemplate2.pptx")